from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLUE_DARK = RGBColor(0x1A, 0x3A, 0x6B)
BLUE_PRIMARY = RGBColor(0x40, 0x9E, 0xFF)
BLUE_LIGHT = RGBColor(0xE6, 0xF0, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x30, 0x30, 0x30)
GRAY_TEXT = RGBColor(0x66, 0x66, 0x66)
RED_CN = RGBColor(0xC0, 0x28, 0x28)
ACCENT_GREEN = RGBColor(0x67, 0xC2, 0x3A)
ACCENT_ORANGE = RGBColor(0xE6, 0xA2, 0x3C)
RED_LIGHT = RGBColor(0xFF, 0xF0, 0xF0)
GOLD = RGBColor(0xC9, 0x9A, 0x2E)

def add_bg(slide, color=BLUE_DARK):
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, color, shape_type=MSO_SHAPE.RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid(); shape.fill.fore_color.rgb = color; shape.line.fill.background()
    return shape

def add_text(slide, left, top, width, height, text, size=14, color=DARK_TEXT, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(size)
    p.font.color.rgb = color; p.font.bold = bold; p.alignment = align
    return txBox

def add_multiline(slide, left, top, width, lines):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.35 * len(lines)))
    tf = txBox.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line.get('text', '')
        p.font.size = Pt(line.get('size', 13))
        p.font.color.rgb = line.get('color', DARK_TEXT)
        p.font.bold = line.get('bold', False)
        p.alignment = line.get('align', PP_ALIGN.LEFT)
    return txBox

# ================================================================
# S1: 封面
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BLUE_DARK)
add_shape(s, Inches(0), Inches(0), Inches(13.333), Inches(0.06), GOLD)
add_shape(s, Inches(0), Inches(7.44), Inches(13.333), Inches(0.06), GOLD)

# 红色竖线装饰
add_shape(s, Inches(4.5), Inches(2.2), Inches(0.06), Inches(1.2), RED_CN)
add_text(s, Inches(5.0), Inches(2.0), Inches(7.5), Inches(0.8),
    "公务用车管理信息系统", size=42, color=WHITE, bold=True)
add_text(s, Inches(5.0), Inches(3.0), Inches(7.5), Inches(0.6),
    "项目建设方案汇报", size=22, color=RGBColor(0xA0, 0xC4, 0xFF))
add_text(s, Inches(5.0), Inches(3.8), Inches(7.5), Inches(0.5),
    "—— 车辆全生命周期数字化管理平台", size=15, color=GRAY_TEXT)

add_shape(s, Inches(5.0), Inches(5.2), Inches(3.5), Inches(0.02), RGBColor(0x55, 0x75, 0xAA))
add_text(s, Inches(5.0), Inches(5.5), Inches(4), Inches(0.4),
    "2026年7月", size=14, color=RGBColor(0x88, 0x99, 0xBB))

# ================================================================
# S2: 建设背景
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
add_shape(s, Inches(0), Inches(0), Inches(13.333), Inches(1.2), BLUE_DARK)
add_text(s, Inches(0.6), Inches(0.25), Inches(12), Inches(0.7),
    "一、建设背景与必要性", size=28, color=WHITE, bold=True)

# 左侧：政策背景
add_shape(s, Inches(0.5), Inches(1.6), Inches(6.0), Inches(2.3), BLUE_LIGHT)
add_text(s, Inches(0.8), Inches(1.7), Inches(5.5), Inches(0.35),
    "政策背景", size=16, color=BLUE_DARK, bold=True)
items1 = [
    "党的十八大以来，中央持续推进公务用车制度改革，要求规范公务用车配备使用管理",
    "《党政机关公务用车管理办法》明确要求建立公务用车管理信息系统，实现全流程监管",
    "国资委要求央企加快推进数字化转型，以信息化手段提升管理效能和合规水平",
    "中央八项规定精神持续深化，对公务用车使用管理提出了更高的纪律要求",
]
y = Inches(2.15)
for it in items1:
    add_text(s, Inches(1.0), y, Inches(5.3), Inches(0.3), "  " + it, size=11, color=DARK_TEXT)
    y += Inches(0.3)

# 右侧：现存问题
add_shape(s, Inches(6.8), Inches(1.6), Inches(6.0), Inches(2.3), RED_LIGHT)
add_text(s, Inches(7.1), Inches(1.7), Inches(5.5), Inches(0.35),
    "当前管理中的突出问题", size=16, color=RED_CN, bold=True)
items2 = [
    "信息孤岛：车辆、驾驶员、费用等数据分散在不同部门，缺乏统一管理平台",
    "审批滞后：纸质审批流程长、效率低，审批记录难以追溯",
    "调度困难：车辆状态不透明，调度靠人工沟通，空置与短缺并存",
    "监管不足：用车行为缺乏有效监督，违规用车难以发现和处置",
]
y = Inches(2.15)
for it in items2:
    add_text(s, Inches(7.3), y, Inches(5.3), Inches(0.3), "  " + it, size=11, color=DARK_TEXT)
    y += Inches(0.3)

# 底部：建设必要性
add_shape(s, Inches(0.5), Inches(4.2), Inches(12.3), Inches(2.8), WHITE)
add_shape(s, Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.04), BLUE_PRIMARY)
add_text(s, Inches(0.8), Inches(4.4), Inches(11.5), Inches(0.35),
    "建设必要性", size=16, color=BLUE_DARK, bold=True)

necessities = [
    ("落实政策要求", "贯彻中央关于公务用车制度改革的决策部署，以信息化手段推动管理规范化、制度化"),
    ("提升管理效能", "实现车辆、驾驶员、调度、费用等信息一体化管理，大幅提升工作效率和管理精细化水平"),
    ("强化合规监管", "建立全流程留痕、可追溯的管理机制，为纪检监察和审计工作提供有力支撑"),
    ("降低运行成本", "通过数据分析优化车辆配置和调度方案，有效降低公务用车运行维护成本"),
]
y = Inches(4.85)
for title, desc in necessities:
    add_shape(s, Inches(0.8), y, Inches(0.05), Inches(0.25), RED_CN)
    add_text(s, Inches(1.1), y - Inches(0.02), Inches(1.8), Inches(0.3), title, size=13, color=RED_CN, bold=True)
    add_text(s, Inches(2.9), y - Inches(0.02), Inches(9.5), Inches(0.3), desc, size=12, color=DARK_TEXT)
    y += Inches(0.4)

# ================================================================
# S3: 系统总体规划
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
add_shape(s, Inches(0), Inches(0), Inches(13.333), Inches(1.2), BLUE_DARK)
add_text(s, Inches(0.6), Inches(0.25), Inches(12), Inches(0.7),
    "二、系统总体规划设计", size=28, color=WHITE, bold=True)

# 整体定位
add_shape(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.2), BLUE_LIGHT)
add_text(s, Inches(0.8), Inches(1.6), Inches(11.5), Inches(0.35),
    "系统定位", size=16, color=BLUE_DARK, bold=True)
add_text(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(0.6),
    '构建覆盖集团总部及下属分公司的统一的公务用车管理信息平台，实现车辆"从购置到报废"的全生命周期管理和用车"从申请到结算"的全业务流程闭环，为管理决策提供数据支撑，为合规监管提供技术保障。',
    size=12, color=DARK_TEXT)

# 四大设计原则
add_text(s, Inches(0.8), Inches(3.0), Inches(11.5), Inches(0.35),
    "设计原则", size=16, color=BLUE_DARK, bold=True)

principles = [
    ("规范统一", "严格遵循公务用车管理相关制度规定，\n建立标准化的业务流程和数据规范", BLUE_PRIMARY),
    ("安全可控", "分级授权、操作留痕、数据加密，\n确保系统安全可靠运行", ACCENT_GREEN),
    ("便捷高效", "优化审批流程、简化操作步骤，\n提升各级人员使用体验和工作效率", ACCENT_ORANGE),
    ("数据驱动", "汇聚车辆运营全维度数据，\n为管理决策提供科学依据", RED_CN),
]

for i, (title, desc, color) in enumerate(principles):
    x = Inches(0.6) + Inches(3.15) * i
    add_shape(s, x, Inches(3.5), Inches(2.9), Inches(1.6), WHITE)
    add_shape(s, x, Inches(3.5), Inches(2.9), Inches(0.06), color)
    add_text(s, x + Inches(0.15), Inches(3.7), Inches(2.6), Inches(0.35), title, size=15, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.15), Inches(4.1), Inches(2.6), Inches(0.8), desc, size=11, color=DARK_TEXT, align=PP_ALIGN.CENTER)

# 组织覆盖范围
add_text(s, Inches(0.8), Inches(5.4), Inches(11.5), Inches(0.35),
    "覆盖范围", size=16, color=BLUE_DARK, bold=True)

scopes = [
    ("组织层级", "集团总部 → 分公司 → 部门\n三级架构全覆盖"),
    ("业务范围", "车辆管理|用车审批|调度派车\n行程记录|费用结算|运营维护"),
    ("角色覆盖", "管理员|员工|部门负责人\n调度员|司机|纪检|财务|分管领导"),
    ("车辆类型", "轿车|SUV|商务车\n中型客车|大型客车"),
]
for i, (title, desc) in enumerate(scopes):
    x = Inches(0.6) + Inches(3.15) * i
    add_shape(s, x, Inches(5.85), Inches(2.9), Inches(1.3), BLUE_LIGHT)
    add_text(s, x + Inches(0.15), Inches(5.95), Inches(2.6), Inches(0.3), title, size=13, color=BLUE_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.15), Inches(6.3), Inches(2.6), Inches(0.7), desc, size=11, color=DARK_TEXT, align=PP_ALIGN.CENTER)

# ================================================================
# S4: 核心业务流程
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
add_shape(s, Inches(0), Inches(0), Inches(13.333), Inches(1.2), BLUE_DARK)
add_text(s, Inches(0.6), Inches(0.25), Inches(12), Inches(0.7),
    "三、核心业务流程设计", size=28, color=WHITE, bold=True)

flow_steps = [
    ("Step 1\n用车申请", "员工线上提交\n选择用车类型\n填写行程信息\n系统自动校验合规", BLUE_PRIMARY),
    ("Step 2\n部门审批", "部门负责人审核\n核对用车合理性\n可选择通过或驳回\n注明审批意见", BLUE_PRIMARY),
    ("Step 3\n领导审批", "分管领导终审\n重点把关合规性\n确认通过或驳回\n审批记录自动存档", BLUE_PRIMARY),
    ("Step 4\n车辆调度", "调度员在线派车\n匹配可用车辆和驾驶员\n确认出车信息\n系统自动更新车辆状态", ACCENT_ORANGE),
    ("Step 5\n出车执行", "驾驶员确认出车\n记录起始里程数\n执行用车任务\n全程行车记录可查", ACCENT_GREEN),
    ("Step 6\n归库登记", "驾驶员确认归库\n记录终止里程数\n系统自动计算里程\n车辆状态恢复可用", ACCENT_GREEN),
    ("Step 7\n费用结算", "录入相关费用\n关联用车申请单\n在线提交报销\n财务审核入账", ACCENT_GREEN),
]

for i, (title, desc, color) in enumerate(flow_steps):
    x = Inches(0.2) + Inches(1.88) * i
    add_shape(s, x, Inches(2.0), Inches(1.7), Inches(2.0), WHITE)
    add_shape(s, x, Inches(2.0), Inches(1.7), Inches(0.06), color)
    add_text(s, x + Inches(0.05), Inches(2.15), Inches(1.6), Inches(0.5), title, size=13, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.08), Inches(2.7), Inches(1.55), Inches(1.2), desc, size=10, color=DARK_TEXT, align=PP_ALIGN.CENTER)

    if i < 6:
        add_text(s, x + Inches(1.65), Inches(2.7), Inches(0.3), Inches(0.3), ">", size=18, color=BLUE_PRIMARY, bold=True)

add_shape(s, Inches(0.3), Inches(4.3), Inches(12.7), Inches(0.02), RGBColor(0xE0, 0xE0, 0xE0))
add_text(s, Inches(0.6), Inches(4.5), Inches(12), Inches(0.35),
    "设计与解决的问题", size=16, color=BLUE_DARK, bold=True)

solved = [
    ("问题一：审批效率低、周期长", "线上审批替代纸质流转，审批人可随时随地处理，审批周期从数天缩短至数小时"),
    ("问题二：车辆调度不透明", "实时展示车辆状态和可用性，调度员一站式完成匹配，消除信息不对称"),
    ("问题三：用车行为难监管", "全流程操作留痕，纪检监察人员可随时查阅任一环节，违规用车无处遁形"),
    ("问题四：费用管理分散混乱", "费用统一关联用车申请，自动汇总统计，预算执行情况一目了然"),
]
y = Inches(4.95)
for title, desc in solved:
    add_text(s, Inches(0.8), y, Inches(5.5), Inches(0.28), title, size=12, color=RED_CN, bold=True)
    add_text(s, Inches(6.5), y, Inches(6.3), Inches(0.28), desc, size=12, color=DARK_TEXT)
    y += Inches(0.35)

# ================================================================
# S5: 功能模块 (上)
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
add_shape(s, Inches(0), Inches(0), Inches(13.333), Inches(1.2), BLUE_DARK)
add_text(s, Inches(0.6), Inches(0.25), Inches(12), Inches(0.7),
    "四、功能模块介绍（上）—— 核心业务模块", size=28, color=WHITE, bold=True)

core_modules = [
    ("车辆档案管理", "建立车辆电子档案，记录购置信息、\n技术参数、保险年检、维修保养等\n全生命周期数据，一车一档", ACCENT_GREEN, "P0"),
    ("驾驶员信息管理", "驾驶员基础档案、驾驶证信息、\n健康状况、培训记录、违规记录等\n统一管理，资质到期自动提醒", ACCENT_GREEN, "P0"),
    ("用车申请管理", "支持6种业务场景的在线申请，\n系统自动校验合规性，关联车辆\n编制配额，提高申请规范性", BLUE_PRIMARY, "P0"),
    ("用车审批管理", "两级审批机制（部门负责人初审\n+ 分管领导终审），支持通过、\n驳回并注明意见，审批留痕", BLUE_PRIMARY, "P0"),
    ("车辆调度管理", "可视化展示车辆和驾驶员可用状态，\n调度员在线完成车辆-驾驶员匹配，\n自动更新车辆状态", BLUE_PRIMARY, "P0"),
    ("出车归库管理", "驾驶员出车前确认车辆状态并记录\n起始里程，归库时记录终止里程，\n系统自动计算行驶里程", BLUE_PRIMARY, "P0"),
    ("费用管理", "统一记录油费、维修费、保险费、\n路桥费等各类费用，支持费用审核\n和报销流程，自动生成费用报表", ACCENT_ORANGE, "P1"),
    ("行程记录查询", "完整记录每次用车行程信息，\n支持按车辆、驾驶员、时间段等\n多维度查询和统计", ACCENT_ORANGE, "P1"),
]

for i, (name, desc, color, priority) in enumerate(core_modules):
    row = i // 4
    col = i % 4
    x = Inches(0.4) + Inches(3.18) * col
    y = Inches(1.6) + Inches(2.85) * row

    add_shape(s, x, y, Inches(2.95), Inches(2.55), WHITE)
    add_shape(s, x, y, Inches(2.95), Inches(0.06), color)
    add_shape(s, x + Inches(2.1), y + Inches(0.1), Inches(0.7), Inches(0.25), color)
    add_text(s, x + Inches(2.1), y + Inches(0.1), Inches(0.7), Inches(0.25), priority, size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.15), y + Inches(0.2), Inches(1.9), Inches(0.3), name, size=14, color=BLUE_DARK, bold=True)
    add_text(s, x + Inches(0.15), y + Inches(0.7), Inches(2.65), Inches(1.7), desc, size=11, color=DARK_TEXT)

# ================================================================
# S6: 功能模块 (下)
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
add_shape(s, Inches(0), Inches(0), Inches(13.333), Inches(1.2), BLUE_DARK)
add_text(s, Inches(0.6), Inches(0.25), Inches(12), Inches(0.7),
    "四、功能模块介绍（下）—— 支撑与监管模块", size=28, color=WHITE, bold=True)

supp_modules = [
    ("车辆状态监控", "实时掌握车辆使用状态和位置，\n空闲/出车/维修/停运一目了然，\n异常状态自动提醒", ACCENT_ORANGE, "P1"),
    ("车辆运营管理", "统一管理加油记录、维修保养、\n保险台账、年检记录、违章处理，\n运营数据一目了然，到期自动提醒", ACCENT_GREEN, "P0"),
    ("驾驶员培训管理", "培训计划制定、培训执行记录、\n考核结果管理，支持8种培训类型，\n培训状态与驾驶资格联动", ACCENT_GREEN, "P0"),
    ("统计分析报表", "提供车辆使用率、费用趋势、\n里程统计等多维度分析报表，\n支持数据导出和图示展示", ACCENT_ORANGE, "P1"),
    ("权限角色管理", "支持管理员、员工、部门负责人、\n调度员、驾驶员、纪检监察、\n财务、分管领导等8种固化角色", RED_CN, "P0"),
    ("通知消息管理", "审批提醒、到期预警、异常告警\n等关键节点消息推送，确保\n重要信息及时传达", ACCENT_ORANGE, "P1"),
]

for i, (name, desc, color, priority) in enumerate(supp_modules):
    row = i // 3
    col = i % 3
    x = Inches(0.4) + Inches(4.22) * col
    y = Inches(1.6) + Inches(2.85) * row

    add_shape(s, x, y, Inches(3.95), Inches(2.55), WHITE)
    add_shape(s, x, y, Inches(3.95), Inches(0.06), color)
    add_shape(s, x + Inches(3.1), y + Inches(0.1), Inches(0.7), Inches(0.25), color)
    add_text(s, x + Inches(3.1), y + Inches(0.1), Inches(0.7), Inches(0.25), priority, size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), y + Inches(0.2), Inches(2.8), Inches(0.3), name, size=14, color=BLUE_DARK, bold=True)
    add_text(s, x + Inches(0.2), y + Inches(0.7), Inches(3.55), Inches(1.7), desc, size=12, color=DARK_TEXT)

# 图例
add_text(s, Inches(0.6), Inches(7.0), Inches(12), Inches(0.3),
    "P0 = 核心功能（优先实现及上线）    P1 = 辅助功能（逐步完善优化）",
    size=11, color=GRAY_TEXT, align=PP_ALIGN.CENTER)

# ================================================================
# S7: 角色权限体系
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
add_shape(s, Inches(0), Inches(0), Inches(13.333), Inches(1.2), BLUE_DARK)
add_text(s, Inches(0.6), Inches(0.25), Inches(12), Inches(0.7),
    "五、角色权限与安全管理", size=28, color=WHITE, bold=True)

add_text(s, Inches(0.6), Inches(1.5), Inches(12), Inches(0.35),
    "8种固化角色，各司其职、权责分明", size=16, color=BLUE_DARK, bold=True)

roles = [
    ("系统管理员", "负责系统配置、用户管理和\n角色权限分配，不参与业务", BLUE_PRIMARY),
    ("普通员工", "提交用车申请、查看个人\n申请记录和审批进度", ACCENT_GREEN),
    ("部门负责人", "审核本部门用车申请，\n管理本部门人员和车辆", ACCENT_ORANGE),
    ("分管领导", "对用车申请进行终审，\n查看分管范围内的全局数据", RED_CN),
    ("车队调度员", "管理车辆和驾驶员调度，\n负责派车、调整和协调", BLUE_PRIMARY),
    ("驾驶员", "执行出车归库确认操作，\n记录行程和里程信息", ACCENT_GREEN),
    ("纪检监察人员", "独立监督用车行为合规性，\n查看操作日志和审计记录", ACCENT_ORANGE),
    ("财务人员", "审核费用报销单据，\n管理预算执行和费用核算", RED_CN),
]

for i, (name, desc, color) in enumerate(roles):
    row = i // 4
    col = i % 4
    x = Inches(0.4) + Inches(3.18) * col
    y = Inches(2.1) + Inches(2.5) * row

    add_shape(s, x, y, Inches(2.95), Inches(2.2), WHITE)
    add_shape(s, x, y, Inches(2.95), Inches(0.06), color)
    # 角色头像圆
    shape = add_shape(s, x + Inches(1.1), y + Inches(0.25), Inches(0.7), Inches(0.7), color, MSO_SHAPE.OVAL)
    add_text(s, x + Inches(1.1), y + Inches(0.4), Inches(0.7), Inches(0.3),
        name[0], size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.15), y + Inches(1.1), Inches(2.65), Inches(0.3),
        name, size=14, color=BLUE_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.15), y + Inches(1.45), Inches(2.65), Inches(0.6),
        desc, size=11, color=GRAY_TEXT, align=PP_ALIGN.CENTER)

# ================================================================
# S8: 解决的实际问题
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
add_shape(s, Inches(0), Inches(0), Inches(13.333), Inches(1.2), BLUE_DARK)
add_text(s, Inches(0.6), Inches(0.25), Inches(12), Inches(0.7),
    "六、系统解决的实际问题", size=28, color=WHITE, bold=True)

problems = [
    ("规范用车行为", "建立标准化线上审批流程，杜绝\n\"先用车后补单\"、私自用车等\n违规现象，从源头强化合规管理",
     "传统模式依赖纸质审批，监管\n滞后；系统实现全流程线上化、\n实时化监管，每笔用车可追溯",
     RED_CN),
    ("提升管理效率", "用车申请、审批、调度全流程线上\n流转，审批人通过移动端随时处理，\n大幅压缩审批周期，减少等待时间",
     "传统审批平均耗时2-3个工作日，\n系统上线后预计缩短至4小时内，\n紧急用车可即时响应",
     BLUE_PRIMARY),
    ("优化资源配置", "实时呈现车辆使用状态和闲置情况，\n为车辆采购和调配提供数据支撑，\n减少不必要的车辆闲置和重复配置",
     "通过运营数据分析，辅助优化\n车辆编制和采购计划，节约\n购置和维护成本",
     ACCENT_GREEN),
    ("强化费用管控", "将油费、维修费、保险费等运营费用\n纳入统一管理平台，预算执行情况\n实时监控，超标自动预警",
     "实现费用支出的全流程追踪，\n防止不合理开支，为年度预算\n编制提供准确数据",
     ACCENT_ORANGE),
    ("保障合规审计", "系统自动记录每个操作节点的时间、\n操作人和操作内容，形成完整的\n审计日志，满足纪检监察和审计要求",
     "纪检人员可调取任意时段、\n任意车辆的完整使用记录，\n实现精准监督",
     RED_CN),
    ("数据辅助决策", "汇聚车辆全生命周期运营数据，\n生成多维度统计分析报表，\n为管理决策提供科学依据",
     "从\"经验管理\"向\"数据管理\"转变，\n用数据说话、用数据决策、\n用数据管理",
     BLUE_PRIMARY),
]

for i, (title, desc, effect, color) in enumerate(problems):
    row = i // 3
    col = i % 3
    x = Inches(0.3) + Inches(4.3) * col
    y = Inches(1.5) + Inches(2.95) * row

    add_shape(s, x, y, Inches(4.05), Inches(2.65), WHITE)
    add_shape(s, x, y, Inches(4.05), Inches(0.06), color)
    add_text(s, x + Inches(0.2), y + Inches(0.2), Inches(3.65), Inches(0.3), title, size=15, color=color, bold=True)
    # 分隔线
    add_shape(s, x + Inches(0.2), y + Inches(0.6), Inches(1.0), Inches(0.03), color)
    add_text(s, x + Inches(0.2), y + Inches(0.7), Inches(1.8), Inches(0.25), "建设效果", size=10, color=color, bold=True)
    add_text(s, x + Inches(2.0), y + Inches(0.7), Inches(1.85), Inches(0.5), effect, size=10, color=DARK_TEXT)
    add_text(s, x + Inches(0.2), y + Inches(1.3), Inches(3.65), Inches(0.25), "功能设计", size=10, color=GRAY_TEXT, bold=True)
    add_text(s, x + Inches(0.2), y + Inches(1.55), Inches(3.65), Inches(0.9), desc, size=11, color=DARK_TEXT)

# ================================================================
# S9: 实施计划
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
add_shape(s, Inches(0), Inches(0), Inches(13.333), Inches(1.2), BLUE_DARK)
add_text(s, Inches(0.6), Inches(0.25), Inches(12), Inches(0.7),
    "七、实施计划与阶段目标", size=28, color=WHITE, bold=True)

phases = [
    ("第一阶段：基础建设（当前阶段）", ACCENT_GREEN,
     ["完成核心业务流程（申请-审批-调度-出车-归库-费用）的系统开发与联调",
      "建立车辆档案、驾驶员档案、部门组织架构等基础数据库",
      "实现8种角色的权限体系，完成用户导入和权限配置",
      "开展系统功能测试和用户验收测试"]),
    ("第二阶段：功能完善", BLUE_PRIMARY,
     ["完善运营管理功能：加油、维修、保险、年检、违章全模块上线",
       "车辆状态实时监控功能部署",
      "统计分析报表和可视化仪表盘开发",
      "系统操作培训和试运行"]),
    ("第三阶段：优化推广", ACCENT_ORANGE,
     ["收集用户反馈，迭代优化系统功能和交互体验",
      "开展全集团推广部署，覆盖所有分公司和部门",
      "探索智能化升级：智能调度、费用预测、异常检测",
      "对接ETC系统、加油卡系统、北斗定位等外部平台"]),
]

for i, (title, color, items) in enumerate(phases):
    x = Inches(0.4) + Inches(4.22) * i
    add_shape(s, x, Inches(1.6), Inches(3.95), Inches(5.5), WHITE)
    add_shape(s, x, Inches(1.6), Inches(3.95), Inches(0.06), color)
    add_shape(s, x, Inches(1.8), Inches(3.95), Inches(0.8), color)
    add_text(s, x + Inches(0.15), Inches(1.9), Inches(3.65), Inches(0.6), title, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    y = Inches(2.85)
    for item in items:
        add_text(s, x + Inches(0.25), y, Inches(3.45), Inches(0.55), "  " + item, size=12, color=DARK_TEXT)
        y += Inches(0.55)
        if items.index(item) < len(items) - 1:
            add_shape(s, x + Inches(0.5), y - Inches(0.08), Inches(2.95), Inches(0.01), RGBColor(0xE8, 0xE8, 0xE8))

# ================================================================
# S10: 预期成效与致谢
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BLUE_DARK)
add_shape(s, Inches(0), Inches(0), Inches(13.333), Inches(0.06), GOLD)
add_shape(s, Inches(0), Inches(7.44), Inches(13.333), Inches(0.06), GOLD)

add_text(s, Inches(0.6), Inches(0.5), Inches(12), Inches(0.7),
    "预期建设成效", size=26, color=WHITE, bold=True)

effects = [
    ("管理规范化", "建立统一的公务用车管理制度\n和标准化操作流程"),
    ("流程高效化", "审批周期从数天缩短至数小时\n紧急用车即时响应"),
    ("监管透明化", "全流程留痕可追溯\n纪检审计有据可查"),
    ("决策数据化", "运营数据自动汇总分析\n管理决策有科学依据"),
    ("成本可控化", "费用预算实时监控\n运行成本显著降低"),
]

for i, (title, desc) in enumerate(effects):
    x = Inches(0.4) + Inches(2.55) * i
    add_shape(s, x, Inches(1.6), Inches(2.3), Inches(2.5), RGBColor(0x22, 0x4B, 0x80))
    add_shape(s, x, Inches(1.6), Inches(2.3), Inches(0.06), GOLD)
    add_text(s, x + Inches(0.1), Inches(2.0), Inches(2.1), Inches(0.4), title, size=16, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.1), Inches(2.6), Inches(2.1), Inches(1.2), desc, size=13, color=RGBColor(0xBB, 0xCC, 0xDD), align=PP_ALIGN.CENTER)

add_shape(s, Inches(4.0), Inches(4.5), Inches(5.3), Inches(0.02), RGBColor(0x55, 0x75, 0xAA))

# 统计数据
stats = [
    ("14", "个功能模块", "全面覆盖公务用车\n管理各业务环节"),
    ("8", "种用户角色", "满足各级各类人员\n使用需求"),
    ("30+", "张数据表", "完整记录车辆全生命\n周期运营数据"),
    ("7x24", "不间断运行", "随时随地处理审批\n提升响应速度"),
]
for i, (num, label, desc) in enumerate(stats):
    x = Inches(0.5) + Inches(3.18) * i
    add_text(s, x, Inches(5.0), Inches(3.0), Inches(0.5), num, size=36, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x, Inches(5.5), Inches(3.0), Inches(0.3), label, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x, Inches(5.85), Inches(3.0), Inches(0.7), desc, size=10, color=RGBColor(0x88, 0x99, 0xBB), align=PP_ALIGN.CENTER)

# 感谢
add_shape(s, Inches(0), Inches(6.7), Inches(13.333), Inches(0.01), RGBColor(0x44, 0x66, 0x99))
add_text(s, Inches(0), Inches(6.85), Inches(13.333), Inches(0.5),
    "恳请各位领导批评指正", size=20, color=RGBColor(0xA0, 0xC4, 0xFF), align=PP_ALIGN.CENTER)

output_path = '/workspace/公务用车管理信息系统-汇报方案.pptx'
prs.save(output_path)
print("Done: " + output_path)
